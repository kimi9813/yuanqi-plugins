import io
import os
import shutil
import tempfile
import zipfile
from pathlib import Path

from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from fastapi.responses import PlainTextResponse
from pydantic import BaseModel, Field

router = APIRouter()

UPLOAD_DIR = Path("data/files")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


class FileReadRequest(BaseModel):
    file_path: str = Field(..., description="文件路径或标识")


class FileConvertResponse(BaseModel):
    filename: str
    mime_type: str
    text: str
    pages: int = 0


def _safe_path(filename: str) -> Path:
    base = UPLOAD_DIR.resolve()
    target = (base / filename).resolve()
    if not str(target).startswith(str(base)):
        raise HTTPException(status_code=400, detail="非法文件路径")
    return target


def _save_upload(upload: UploadFile) -> Path:
    target = _safe_path(upload.filename or "unnamed")
    with target.open("wb") as f:
        shutil.copyfileobj(upload.file, f)
    return target


def _read_text_file(path: Path) -> str:
    encodings = ["utf-8", "gbk", "latin-1"]
    for enc in encodings:
        try:
            return path.read_text(encoding=enc)
        except UnicodeDecodeError:
            continue
    raise HTTPException(status_code=400, detail="无法解码文本文件")


def _read_pdf(path: Path) -> tuple:
    try:
        import PyPDF2

        reader = PyPDF2.PdfReader(str(path))
        pages = len(reader.pages)
        text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
        return text, pages
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"PDF 解析失败: {exc}")


def _read_docx(path: Path) -> tuple:
    try:
        import docx

        doc = docx.Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text)
        return text, len(doc.sections)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Word 解析失败: {exc}")


def _read_excel(path: Path) -> tuple:
    try:
        import openpyxl

        wb = openpyxl.load_workbook(path, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                rows.append("\t".join(str(cell) if cell is not None else "" for cell in row))
            parts.append(f"--- Sheet: {sheet.title} ---\n" + "\n".join(rows))
        return "\n\n".join(parts), len(wb.worksheets)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Excel 解析失败: {exc}")


def _read_pptx(path: Path) -> tuple:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides), len(prs.slides)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"PPT 解析失败: {exc}")


def _read_zip(path: Path) -> str:
    try:
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            return "\n".join(names[:200])
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"ZIP 读取失败: {exc}")


def _convert_file(path: Path) -> FileConvertResponse:
    suffix = path.suffix.lower()
    mime = {
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".py": "text/x-python",
        ".sh": "text/x-shellscript",
        ".js": "text/javascript",
        ".json": "application/json",
        ".csv": "text/csv",
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".zip": "application/zip",
    }.get(suffix, "application/octet-stream")

    text, pages = "", 0
    if suffix in [".txt", ".md", ".py", ".sh", ".js", ".json", ".csv"]:
        text = _read_text_file(path)
    elif suffix == ".pdf":
        text, pages = _read_pdf(path)
    elif suffix == ".docx":
        text, pages = _read_docx(path)
    elif suffix == ".xlsx":
        text, pages = _read_excel(path)
    elif suffix == ".pptx":
        text, pages = _read_pptx(path)
    elif suffix == ".zip":
        text = _read_zip(path)
    else:
        text = f"[binary file: {path.name}]"

    return FileConvertResponse(filename=path.name, mime_type=mime, text=text[:20000], pages=pages)


@router.post("/read", operation_id="file_read")
async def file_read(file_path: str = Form(..., description="相对上传目录的文件路径")):
    """读取并转换已上传文件为文本"""
    path = _safe_path(file_path)
    if not path.exists():
        raise HTTPException(status_code=404, detail="文件不存在")
    return _convert_file(path)


@router.post("/upload-and-convert", operation_id="file_upload_and_convert")
async def file_upload_and_convert(
    file: UploadFile = File(..., description="待上传并转换的文件"),
):
    """上传文件并立即转换为文本"""
    try:
        path = _save_upload(file)
        return _convert_file(path)
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"上传处理失败: {exc}")


@router.get("/list", operation_id="file_list")
async def file_list():
    """列出已上传文件"""
    files = []
    for p in UPLOAD_DIR.iterdir():
        if p.is_file():
            files.append({"name": p.name, "size": p.stat().st_size})
    return {"files": files}
